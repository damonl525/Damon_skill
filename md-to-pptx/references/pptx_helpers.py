#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX Helper Functions Template
===============================
Reusable building blocks for generating professional PPTX presentations.
Designed for technical/statistical content with formulas, tables, code blocks,
and callout boxes.

Usage:
  1. Copy this file or import its functions into your gen_pptx.py
  2. Customize slide content using the helpers below
  3. Run: python gen_pptx.py

Dependencies: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def _emu(v):
    """Convert bare number to EMU via Inches(); pass through pre-converted values.

    python-pptx uses EMU units (914400 per inch). Inches() returns an int (Emu
    subclass), so isinstance(v, int) matches pre-converted values too. We use a
    heuristic: ints > 100 are already EMU (since slide coordinates are < 14).
    """
    if isinstance(v, int):
        return v if v > 100 else Inches(v)
    if isinstance(v, float):
        return Inches(v)
    return v

# ══════════════════════════════════════════════════════════════════════════════
# STYLE CONSTANTS — Customize these to change the visual theme
# ══════════════════════════════════════════════════════════════════════════════

# Fonts
FN = '微软雅黑'           # All non-code text
FN_CODE = 'Consolas'      # Code blocks and formulas

# Font Sizes (pt)
SZ_TITLE = 38             # Slide big title
SZ_SECTION = 32           # Section heading
SZ_SUBTITLE = 22          # Subtitle
SZ_H3 = 20                # Sub-heading (h3)
SZ_BODY = 17              # Body text
SZ_BODY_SM = 16           # Body small
SZ_LIST = 17              # List items
SZ_FORMULA = 16           # Formula text (Consolas)
SZ_CODE = 13              # Code text (Consolas)
SZ_CALLOUT_TITLE = 16     # Callout bold title
SZ_CALLOUT_BODY = 15      # Callout body
SZ_TABLE_HDR = 15         # Table header (white on green)
SZ_TABLE_CELL = 14        # Table cell
SZ_FLOW = 14              # Flow diagram boxes
SZ_SLIDE_NUM = 11         # Slide number
SZ_LABEL = 17             # Labels

# Colors
ACCENT      = RGBColor(0x2D, 0x6A, 0x4F)   # Dark green
ACCENT_LIGHT= RGBColor(0x40, 0x91, 0x6C)   # Medium green
ACCENT_BG   = RGBColor(0xD8, 0xF3, 0xDC)   # Light green background
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x2E)   # Near-black body text
GRAY        = RGBColor(0x6C, 0x75, 0x7D)   # Subtitles, slide numbers
RED         = RGBColor(0xD0, 0x00, 0x00)   # Warning accent
RED_BG      = RGBColor(0xFC, 0xE4, 0xE4)   # Warning background
CODE_BG     = RGBColor(0xF0, 0xF0, 0xF0)   # Code block background
BORDER      = RGBColor(0xDE, 0xE2, 0xE6)   # Borders

# Alignment shortcuts (AI models often use these instead of PP_ALIGN.XXX)
LEFT   = PP_ALIGN.LEFT
CENTER = PP_ALIGN.CENTER
RIGHT  = PP_ALIGN.RIGHT
DARK_GREEN  = RGBColor(0x1B, 0x43, 0x32)   # Formula text
DARK_RED    = RGBColor(0x64, 0x12, 0x20)   # Warning text

# Slide Dimensions (16:9 widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Layout Constants
LEFT_BAR_W = Inches(0.06)   # Accent bar width
CONTENT_LEFT = Inches(0.5)  # Main content left margin
CONTENT_WIDTH = Inches(12.3) # Main content usable width
COL_W = Inches(5.8)         # Single column width
COL_L_LEFT = Inches(0.5)    # Left column start
COL_R_LEFT = Inches(6.8)    # Right column start


# ══════════════════════════════════════════════════════════════════════════════
# PRESENTATION SETUP
# ══════════════════════════════════════════════════════════════════════════════

def create_presentation():
    """Create a new widescreen presentation and return (prs, blank_layout)."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]
    return prs, blank_layout


# ══════════════════════════════════════════════════════════════════════════════
# BASIC HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def add_slide(prs, blank_layout):
    """Add a new slide with left accent bar. Returns the slide object."""
    slide = prs.slides.add_slide(blank_layout)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), LEFT_BAR_W, SLIDE_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    return slide


def add_tb(slide, left, top, width, height):
    """Add a textbox at the specified position."""
    return slide.shapes.add_textbox(_emu(left), _emu(top), _emu(width), _emu(height))


def set_text(tf, text, sz=SZ_BODY, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    """Set text in a text frame (replaces existing content)."""
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FN
    return r


def add_para(tf, text, sz=SZ_BODY, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, sb=Pt(5)):
    """Append a new paragraph to a text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = sb
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FN
    return r


def add_bg(slide, left, top, w, h, fill, border_color=None):
    """Add a rounded rectangle background shape."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, _emu(left), _emu(top), _emu(w), _emu(h)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


# ══════════════════════════════════════════════════════════════════════════════
# COMPOSITE ELEMENTS
# ══════════════════════════════════════════════════════════════════════════════

def add_slide_num(slide, num, total):
    """Add slide number (bottom-right corner)."""
    tb = add_tb(slide, Inches(12.4), Inches(7.0), Inches(0.8), Inches(0.4))
    set_text(tb.text_frame, f'{num} / {total}',
             sz=SZ_SLIDE_NUM, color=GRAY, align=PP_ALIGN.RIGHT)


def add_title(slide, text, top=Inches(0.4)):
    """Add large slide title."""
    tb = add_tb(slide, CONTENT_LEFT, top, CONTENT_WIDTH, Inches(0.8))
    set_text(tb.text_frame, text, sz=SZ_TITLE, bold=True, color=ACCENT)


def add_subtitle(slide, text, top=Inches(1.1)):
    """Add subtitle text."""
    tb = add_tb(slide, CONTENT_LEFT, top, CONTENT_WIDTH, Inches(0.6))
    set_text(tb.text_frame, text, sz=SZ_SUBTITLE, color=GRAY)


def add_section(slide, text, top=Inches(0.25)):
    """Add section heading with green underline bar."""
    tb = add_tb(slide, CONTENT_LEFT, top, CONTENT_WIDTH, Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(SZ_SECTION)
    r.font.bold = True
    r.font.color.rgb = BLACK
    r.font.name = FN
    # Green underline bar
    ln = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        CONTENT_LEFT, top + Inches(0.72),
        Inches(3), Inches(0.04)
    )
    ln.fill.solid()
    ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background()


def add_h3(slide, text, top=Inches(0.3)):
    """Add H3 sub-heading (bold, no underline bar)."""
    tb = add_tb(slide, CONTENT_LEFT, top, CONTENT_WIDTH, Inches(0.45))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(SZ_H3)
    r.font.bold = True
    r.font.color.rgb = BLACK
    r.font.name = FN


def add_formula(slide, text, left, top, width, height):
    """
    Add a formula box with green background and left accent bar.

    text: multi-line string (each line becomes a paragraph)
    """
    left, top, width, height = _emu(left), _emu(top), _emu(width), _emu(height)
    add_bg(slide, left, top, width, height, ACCENT_BG, ACCENT_LIGHT)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(0.05), height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    tb = add_tb(slide, left + Inches(0.25), top + Inches(0.12),
                width - Inches(0.4), height - Inches(0.18))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.strip().split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(SZ_FORMULA)
        r.font.name = FN_CODE
        r.font.color.rgb = DARK_GREEN


def add_callout(slide, title, body, left, top, width, height, warn=False):
    """
    Add a callout box (info or warning style).

    title: bold heading text
    body: multi-line string body text
    warn: True for red warning style, False for green info style
    """
    left, top, width, height = _emu(left), _emu(top), _emu(width), _emu(height)
    if warn:
        bg_c, bd_c, tx_c = RED_BG, RED, DARK_RED
    else:
        bg_c, bd_c, tx_c = ACCENT_BG, ACCENT_LIGHT, DARK_GREEN

    add_bg(slide, left, top, width, height, bg_c, bd_c)
    tb = add_tb(slide, left + Inches(0.2), top + Inches(0.1),
                width - Inches(0.4), height - Inches(0.18))
    tf = tb.text_frame
    tf.word_wrap = True

    # Title
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(SZ_CALLOUT_TITLE)
    r.font.bold = True
    r.font.color.rgb = tx_c
    r.font.name = FN

    # Body lines
    for line in body.strip().split('\n'):
        p2 = tf.add_paragraph()
        p2.space_before = Pt(3)
        r2 = p2.add_run()
        r2.text = line
        r2.font.size = Pt(SZ_CALLOUT_BODY)
        r2.font.color.rgb = tx_c
        r2.font.name = FN


def add_code(slide, text, left, top, width, height):
    """
    Add a code block with gray background.

    text: multi-line string code content
    """
    left, top, width, height = _emu(left), _emu(top), _emu(width), _emu(height)
    add_bg(slide, left, top, width, height, CODE_BG, BORDER)
    tb = add_tb(slide, left + Inches(0.2), top + Inches(0.1),
                width - Inches(0.4), height - Inches(0.18))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.strip().split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(SZ_CODE)
        r.font.name = FN_CODE
        r.font.color.rgb = BLACK


def add_table(slide, headers, rows, left, top, width, height):
    """
    Add a formatted table with green header row.

    headers: list of column header strings
    rows: list of lists (each inner list = one data row)
    """
    if not headers:
        raise ValueError("headers cannot be empty")
    if not rows:
        raise ValueError("rows cannot be empty — use add_bg for empty placeholders")
    n_rows = len(rows) + 1
    n_cols = len(headers)
    left, top, width, height = _emu(left), _emu(top), _emu(width), _emu(height)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table
    col_w = int(width / n_cols)
    for j in range(n_cols):
        tbl.columns[j].width = col_w

    # Header row
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(SZ_TABLE_HDR)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FN
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i + 1, j)
            c.text = str(val)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(SZ_TABLE_CELL)
                p.font.color.rgb = BLACK
                p.font.name = FN
            if i % 2 == 1:
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)


def add_flow_boxes(slide, items, top, box_height=Inches(0.9)):
    """
    Add horizontal flow boxes with arrows between them.

    items: list of (text, left_position) tuples
    top: vertical position for the row
    """
    if not items:
        raise ValueError("items cannot be empty")
    top, box_height = _emu(top), _emu(box_height)
    box_w = Inches(2.6)
    for text, left in items:
        add_bg(slide, left, top, box_w, box_height, ACCENT_BG, ACCENT_LIGHT)
        tb = add_tb(slide, left + Inches(0.1), top + Inches(0.05),
                    box_w - Inches(0.2), box_height - Inches(0.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.size = Pt(SZ_FLOW)
        r.font.bold = True
        r.font.color.rgb = ACCENT
        r.font.name = FN

    # Arrows between boxes
    for i in range(len(items) - 1):
        _, left_current = items[i]
        arrow_left = left_current + box_w + Inches(0.05)
        a = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            arrow_left, top + Inches(0.22),
            Inches(0.35), Inches(0.35)
        )
        a.fill.solid()
        a.fill.fore_color.rgb = ACCENT
        a.line.fill.background()


def add_comparison(slide, left_title, left_items, right_title, right_items,
                   left_pos=COL_L_LEFT, right_pos=COL_R_LEFT,
                   top=Inches(4.9), width=COL_W, height=Inches(2.0),
                   left_warn=True):
    """
    Add side-by-side comparison boxes (red vs green).

    left_title: left box heading
    left_items: list of bullet point strings for left box
    right_title: right box heading
    right_items: list of bullet point strings for right box
    left_warn: True = left box is warning (red), False = left box is info (green)
    """
    left_pos, right_pos, top, width, height = (
        _emu(left_pos), _emu(right_pos), _emu(top), _emu(width), _emu(height)
    )
    # Left box
    left_bg = RED_BG if left_warn else ACCENT_BG
    left_bd = RED if left_warn else ACCENT_LIGHT
    left_tx = DARK_RED if left_warn else DARK_GREEN
    left_hd = RED if left_warn else ACCENT
    add_bg(slide, left_pos, top, width, height, left_bg, left_bd)
    tb_l = add_tb(slide, left_pos + Inches(0.2), top + Inches(0.1),
                  width - Inches(0.4), height - Inches(0.18))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    r = tf_l.paragraphs[0].add_run()
    r.text = left_title
    r.font.size = Pt(SZ_H3)
    r.font.bold = True
    r.font.color.rgb = left_hd
    r.font.name = FN
    for item in left_items:
        p = tf_l.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = '• ' + item
        r.font.size = Pt(SZ_BODY_SM)
        r.font.color.rgb = left_tx
        r.font.name = FN

    # Right box
    right_bg = ACCENT_BG if left_warn else RED_BG
    right_bd = ACCENT_LIGHT if left_warn else RED
    right_tx = DARK_GREEN if left_warn else DARK_RED
    right_hd = ACCENT if left_warn else RED
    add_bg(slide, right_pos, top, width, height, right_bg, right_bd)
    tb_r = add_tb(slide, right_pos + Inches(0.2), top + Inches(0.1),
                  width - Inches(0.4), height - Inches(0.18))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    r = tf_r.paragraphs[0].add_run()
    r.text = right_title
    r.font.size = Pt(SZ_H3)
    r.font.bold = True
    r.font.color.rgb = right_hd
    r.font.name = FN
    for item in right_items:
        p = tf_r.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = '• ' + item
        r.font.size = Pt(SZ_BODY_SM)
        r.font.color.rgb = right_tx
        r.font.name = FN

    # "vs." label
    vs_left = left_pos + width + Inches(0.05)
    tb_vs = add_tb(slide, vs_left, top + Inches(0.5), Inches(0.6), Inches(0.5))
    set_text(tb_vs.text_frame, 'vs.', sz=24, bold=True,
             color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# BULLET LIST HELPER
# ══════════════════════════════════════════════════════════════════════════════

def add_bullet_list(slide, items, left, top, width, height,
                    sz=SZ_LIST, bold=False, color=BLACK, bullet='• '):
    """Add a simple bullet list."""
    if not items:
        raise ValueError("items cannot be empty")
    tb = add_tb(slide, _emu(left), _emu(top), _emu(width), _emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(6)
        r = p.add_run()
        r.text = bullet + item
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FN


# ══════════════════════════════════════════════════════════════════════════════
# AUTO-LAYOUT (Recommended — no manual coordinates needed)
# ══════════════════════════════════════════════════════════════════════════════

class Layout:
    """Auto-layout manager: stack elements vertically without coordinates.

    Tracks a Y cursor and auto-positions each element below the previous one.
    Height is estimated from content. Returns self for chaining.

    Usage:
        prs, layout = create_presentation()
        total = 8

        s = add_slide(prs, layout)
        lm = Layout(s)
        lm.title("Score Method for Missing Data")
        lm.subtitle("Rubin's Rules & Variance Estimation")
        lm.bullets(["Obj 1", "Obj 2", "Obj 3"])
        lm.slide_num(1, total)

        prs.save("output.pptx")
    """

    _TOP = 0.2
    _BOTTOM = 0.45
    _GAP = 0.15

    def __init__(self, slide):
        self._slide = slide
        self._y = self._TOP

    @property
    def remaining(self):
        """Remaining usable height (inches) before bottom margin."""
        return 7.5 - self._y - self._BOTTOM

    @property
    def cursor(self):
        """Current Y position (inches from top)."""
        return self._y

    def reset(self):
        """Reset cursor to top (for reusing Layout on a new slide)."""
        self._y = self._TOP
        return self

    def _advance(self, height):
        y = self._y
        self._y += height + self._GAP
        return y

    def gap(self, inches=0.2):
        """Add explicit vertical gap."""
        self._y += inches
        return self

    def advance(self, inches):
        """Manually advance cursor (for mixing auto + manual layouts)."""
        self._y += inches
        return self

    # ── Slide-level ──

    def slide_num(self, num, total):
        """Add slide number (bottom-right, no cursor advance)."""
        add_slide_num(self._slide, num, total)
        return self

    # ── Text ──

    def title(self, text):
        """Add slide title (large, green, bold)."""
        y = self._advance(0.8)
        add_title(self._slide, text, top=_emu(y))
        return self

    def subtitle(self, text):
        """Add subtitle (gray, smaller)."""
        y = self._advance(0.6)
        add_subtitle(self._slide, text, top=_emu(y))
        return self

    def section(self, text):
        """Add section heading with green underline bar."""
        y = self._advance(0.85)
        add_section(self._slide, text, top=_emu(y))
        return self

    def h3(self, text):
        """Add H3 sub-heading (bold, no underline)."""
        y = self._advance(0.5)
        add_h3(self._slide, text, top=_emu(y))
        return self

    # ── Content (full-width, auto-height) ──

    def bullets(self, items, sz=SZ_LIST, bold=False):
        """Add bullet list. Height auto-estimated from item count."""
        h = max(0.5, len(items) * 0.38 + 0.1)
        y = self._advance(h)
        add_bullet_list(self._slide, items, CONTENT_LEFT, _emu(y),
                        CONTENT_WIDTH, _emu(h), sz=sz, bold=bold)
        return self

    def formula(self, text):
        """Add formula box. Height auto-estimated from line count."""
        lines = text.strip().count('\n') + 1
        h = max(0.6, lines * 0.35 + 0.3)
        y = self._advance(h)
        add_formula(self._slide, text, CONTENT_LEFT, _emu(y),
                    CONTENT_WIDTH, _emu(h))
        return self

    def code(self, text):
        """Add code block. Height auto-estimated from line count."""
        lines = text.strip().count('\n') + 1
        h = max(0.5, lines * 0.25 + 0.2)
        y = self._advance(h)
        add_code(self._slide, text, CONTENT_LEFT, _emu(y),
                 CONTENT_WIDTH, _emu(h))
        return self

    def callout(self, title, body, warn=False):
        """Add callout box (info or warning). Height auto-estimated."""
        lines = body.strip().count('\n') + 1
        h = max(0.7, lines * 0.3 + 0.55)
        y = self._advance(h)
        add_callout(self._slide, title, body, CONTENT_LEFT, _emu(y),
                    CONTENT_WIDTH, _emu(h), warn=warn)
        return self

    def table(self, headers, rows):
        """Add table. Height auto-estimated from row count."""
        h = max(0.7, (len(rows) + 1) * 0.38 + 0.1)
        y = self._advance(h)
        add_table(self._slide, headers, rows, CONTENT_LEFT, _emu(y),
                  CONTENT_WIDTH, _emu(h))
        return self

    # ── Composite (multi-element patterns) ──

    def comparison(self, left_title, left_items, right_title, right_items,
                   left_warn=True):
        """Add side-by-side comparison boxes (red vs green)."""
        h = 2.0
        y = self._advance(h)
        add_comparison(self._slide, left_title, left_items,
                       right_title, right_items, top=_emu(y),
                       left_warn=left_warn)
        return self

    def flow(self, items, box_height=0.9):
        """Add flow diagram. items: list of text strings (auto-positioned)."""
        y = self._advance(box_height)
        n = len(items)
        if n == 0:
            return self
        box_w_in = 2.6
        arrow_in = 0.4
        total_needed = n * box_w_in + (n - 1) * arrow_in
        start_x = 0.5 + (12.3 - total_needed) / 2
        positioned = []
        for i, text in enumerate(items):
            x = start_x + i * (box_w_in + arrow_in)
            positioned.append((text, _emu(x)))
        add_flow_boxes(self._slide, positioned, top=_emu(y),
                       box_height=_emu(box_height))
        return self

    # ── Two-column layouts ──

    def two_col_bullets(self, left_title, left_items, right_title, right_items):
        """Two columns: h3 title + bullet list in each column."""
        max_n = max(len(left_items), len(right_items))
        h = max(1.5, max_n * 0.38 + 0.65)
        y = self._advance(h)
        # Left
        tb_l = add_tb(self._slide, COL_L_LEFT, _emu(y), COL_W, _emu(0.5))
        set_text(tb_l.text_frame, left_title, sz=SZ_H3, bold=True)
        add_bullet_list(self._slide, left_items, COL_L_LEFT, _emu(y + 0.55),
                        COL_W, _emu(h - 0.55))
        # Right
        tb_r = add_tb(self._slide, COL_R_LEFT, _emu(y), COL_W, _emu(0.5))
        set_text(tb_r.text_frame, right_title, sz=SZ_H3, bold=True)
        add_bullet_list(self._slide, right_items, COL_R_LEFT, _emu(y + 0.55),
                        COL_W, _emu(h - 0.55))
        return self

    def two_col_table_code(self, table_headers, table_rows, code_text):
        """Two columns: table on left, code block on right."""
        tbl_h = max(0.7, (len(table_rows) + 1) * 0.38 + 0.1)
        code_lines = code_text.strip().count('\n') + 1
        code_h = max(0.5, code_lines * 0.25 + 0.2)
        h = max(tbl_h, code_h)
        y = self._advance(h)
        add_table(self._slide, table_headers, table_rows,
                  COL_L_LEFT, _emu(y), COL_W, _emu(tbl_h))
        add_code(self._slide, code_text,
                 COL_R_LEFT, _emu(y), COL_W, _emu(code_h))
        return self

    def two_col_bullets_table(self, bullet_title, items, table_headers, table_rows):
        """Two columns: bullet list on left, table on right."""
        blt_h = max(0.5, len(items) * 0.38 + 0.1)
        tbl_h = max(0.7, (len(table_rows) + 1) * 0.38 + 0.1)
        h = max(blt_h + 0.5, tbl_h)
        y = self._advance(h)
        # Left: h3 + bullets
        tb_l = add_tb(self._slide, COL_L_LEFT, _emu(y), COL_W, _emu(0.5))
        set_text(tb_l.text_frame, bullet_title, sz=SZ_H3, bold=True)
        add_bullet_list(self._slide, items, COL_L_LEFT, _emu(y + 0.55),
                        COL_W, _emu(h - 0.55))
        # Right: table
        add_table(self._slide, table_headers, table_rows,
                  COL_R_LEFT, _emu(y), COL_W, _emu(tbl_h))
        return self


# ══════════════════════════════════════════════════════════════════════════════
# EXAMPLE USAGE
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == '__main__':
    import os

    prs, layout = create_presentation()
    total = 6

    # ── Slide 1: Title (auto-layout, no coordinates) ──
    s = add_slide(prs, layout)
    Layout(s).title('Presentation Title').subtitle('Subtitle text here').slide_num(1, total)

    # ── Slide 2: Section + content ──
    s = add_slide(prs, layout)
    lm = Layout(s)
    lm.section('Section Title')
    lm.formula('Z = (d - delta) / sqrt(V)')
    lm.callout('Key Insight', 'This is an important callout box.')
    lm.slide_num(2, total)

    # ── Slide 3: Two-column bullets + table ──
    s = add_slide(prs, layout)
    lm = Layout(s)
    lm.h3('Clinical Trial Design')
    lm.two_col_bullets_table(
        'Key Features',
        ['Phase III randomized trial', 'Primary endpoint: PFS', 'n = 350 per arm'],
        ['Factor', 'Value'],
        [['Design', '1:1 randomization'], ['Population', 'HER2-negative']],
    )
    lm.slide_num(3, total)

    # ── Slide 4: Two-column table + code ──
    s = add_slide(prs, layout)
    lm = Layout(s)
    lm.h3('Implementation')
    lm.two_col_table_code(
        ['Method', 'Variance', 'CI Type'],
        [['Wald-Rubin', 'Rubin', 'Wald'],
         ['Score', 'Rubin', 'Score inversion']],
        'score_test <- function(d, V) {\n  Z <- d / sqrt(V)\n  2 * pnorm(-abs(Z))\n}',
    )
    lm.slide_num(4, total)

    # ── Slide 5: Flow + Comparison ──
    s = add_slide(prs, layout)
    lm = Layout(s)
    lm.h3('Workflow')
    lm.flow(['Collect Data', 'MI (m=20)', 'Pool', 'Score Test'])
    lm.comparison('Wald', ['Simple', 'May undercover'], 'Score', ['Precise', 'Better coverage'])
    lm.slide_num(5, total)

    # ── Slide 6: Summary ──
    s = add_slide(prs, layout)
    Layout(s).section('Summary').bullets([
        'Auto-layout eliminates manual positioning',
        'Height auto-estimated from content',
        'Two-column helpers for common patterns',
    ]).slide_num(6, total)

    # Save
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'example_output.pptx')
    prs.save(out)
    print(f'Saved: {out}')
